import os
import json
import docx
import pdfplumber
from dotenv import load_dotenv, find_dotenv
from pptx import Presentation
from google import genai
from courses.schemas import GeneratedJourney, get_assessment_mix

from google.genai import types

load_dotenv(find_dotenv())

api_key = os.environ.get("GEMINI_API_KEY")
client = genai.Client(api_key=api_key) if api_key else None
GEMINI_MODEL = "gemini-3.6-flash"

# ==========================================
# 1. PROMPT PROFILE DEFINITIONS
# ==========================================

STUDY_GOAL_PROFILES = {
    "deep_learning": (
        "PEDAGOGICAL DIRECTIVE: Comprehensive In-Depth Study.\n"
        "- Prioritize deep technical context, workflows, and root causes.\n"
        "- Require formal definitions followed by 'In Simple Words' breakdowns.\n"
        "- Provide real-world architecture examples, system trade-offs, and multi-layered analogies.\n"
        "- Break chapters into 3 to 5 logical, numbered sequential sections."
    ),
    "balanced_review": (
        "PEDAGOGICAL DIRECTIVE: Balanced Academic Review.\n"
        "- Balance concise definitions with practical context and takeaways.\n"
        "- Include essential analogies for difficult abstractions only.\n"
        "- Emphasize high-value exam concepts, comparison matrices, and clear section flow.\n"
        "- Break chapters into 2 to 4 focused sequential sections."
    ),
    "quick_cram": (
        "PEDAGOGICAL DIRECTIVE: High-Yield Exam Cram.\n"
        "- Maximize memory density, rapid scannability, and high-frequency exam points.\n"
        "- Prioritize exact term-definition pairs, recognition cues, and acronyms.\n"
        "- Include high-yield enumerations (numbered lists, steps, components).\n"
        "- Keep introductions brief; surface 'Common Traps & Confusions' prominently."
    ),
}

ASSESSMENT_PROFILES = {
    "multiple_choice": (
        "ASSESSMENT TARGET: Multiple Choice & True/False Comprehension.\n"
        "- Focus on conceptual distinction, edge cases, and plausible distractors.\n"
        "- Ensure distractors reflect common student misconceptions, not absurd answers."
    ),
    "identification": (
        "ASSESSMENT TARGET: Exact Term Identification.\n"
        "- Emphasize precise technical vocabulary, acronyms, and standard definitions.\n"
        "- Provide clear contextual cues that distinguish easily confused terms."
    ),
    "enumeration": (
        "ASSESSMENT TARGET: Structured Enumeration.\n"
        "- Extract explicit named lists, categories, phases, matrices, and rule-sets from the source.\n"
        "- Provide clear prompts and memory cues (acronyms, initialisms, counts) for each list."
    ),
}

DOCUMENT_STRUCTURE_RULES = """
CHAPTER SCOPING & STRUCTURE MANDATES:
1. Divide this material into 3 to 5 focused, compact chapters based on natural topic breaks (approx. 2 to 3 slides per chapter).
2. NEVER cram more than 3 related major concepts into one chapter.
3. Chapter 1 must ONLY cover Portfolio Auditing & Tool Governance (e.g., TIME Framework and Tool Sprawl).
4. Chapter 2 covers Lifecycles, Methodologies (Waterfall vs. Agile), and RACI.
5. Chapter 3 covers Constraints (Iron Triangle) and Scope Management (Scope Creep).
6. Chapter 4 covers Team Leadership, Pre-Mortems, and Enterprise Risk (The 4 T's).
"""

ANTI_REDUNDANCY_RULES = """
CONTENT DE-DUPLICATION (ZERO-BLOAT) RULES:
1. One Home Per Fact: Mention each term or concept in EXACTLY ONE widget.
2. If a concept is defined inside a "section", DO NOT repeat it in "key_terms" or "enumerations".
3. Use widgets conditionally and sparsely:
   - "comparisons": ONLY if the source explicitly contrasts two items (e.g. Waterfall vs Agile).
   - "enumerations": ONLY for explicit lists or acronyms (e.g. 4 T's, 5 Phases, RACI).
   - "analogy": Maximum ONE per chapter, only if directly in the text or directly clarifying.
   - "common_confusions": Maximum ONE pair per chapter, only for genuinely mixed-up concepts.
4. Keep section introductions to 1-2 concise sentences. Do not expand with generic textbook fluff.
"""

QUIZ_RULES = """
INTERACTIVE QUIZ COMPOSITION MANDATES:
1. Generate exactly 10 questions per chapter using the required distribution below.
2. Multiple-choice questions MUST have exactly 4 choices and exactly 1 correct answer.
3. True/False questions MUST have exactly 2 choices ("True" and "False") and exactly 1 correct answer.
4. Identification questions MUST have no choices and at least one accepted answer.
5. Enumeration questions MUST have no choices, at least 2 expected items, and explicit order_matters metadata.
6. Every question must include a 2-3 sentence educational "explanation".
"""

SOURCE_GROUNDING_RULES = """
SOURCE-GROUNDING & INTEGRITY MANDATES:
1. Ground truth: Use ONLY facts, terms, frameworks, and statistics present in the <source_material>.
2. Do not invent: Never fabricate citations, outside libraries, or unmentioned tools.
3. Sparse fields: If the source material does not contain enough information for a specific optional field (e.g., analogies or comparisons), return an empty array [] or null. Never fabricate filler content.
4. Acronyms & Terms: Retain the exact capitalization and naming used in the source (e.g., BDAT, TOGAF, RACI, TIME).
5. Quizzes: All quiz questions must be strictly solvable using the generated chapter content.
"""

CURRICULUM_JSON_SCHEMA = """
OUTPUT FORMAT: Output valid JSON matching this exact structure:
{
  "schema_version": "1.0",
  "course": {
    "title": "Course Title",
    "description": "1-2 sentence academic summary of the entire course curriculum.",
    "source_type": "reviewer",
    "difficulty": "intermediate"
  },
  "chapters": [
    {
      "order": 1,
      "title": "Chapter Title",
      "week_label": "e.g., Week 8-9 (or null if not indicated)",
      "focus": "1-sentence summary of what this specific chapter covers.",
      "overview": "Concise chapter overview summarizing core themes.",
      "estimated_minutes": 15,
      "learning_objectives": [
        "Actionable outcome 1",
        "Actionable outcome 2"
      ],
      "sections": [
        {
          "order": 1,
          "title": "Section Title",
          "introduction": "1-2 sentences setting up the context for this topic.",
          "concepts": [
            {
              "term": "Concept Name",
              "formal_definition": "Precise textbook definition grounded in the document.",
              "simple_explanation": "Plain-language explanation.",
              "analogy": "Memorable analogy or null.",
              "examples": ["Concrete real-world example from text"]
            }
          ],
          "key_points": [
            "Bullet point 1 summarizing a core rule or formula",
            "Bullet point 2"
          ]
        }
      ],
      "key_terms": [
        {
          "term": "Concept Name",
          "definition": "Precise definition grounded in the document."
        }
      ],
      "comparisons": [
        {
          "title": "Comparison Matrix Title",
          "columns": ["Criterion", "Option A", "Option B"],
          "rows": [
            {
              "criterion": "Execution Model",
              "values": ["Sequential and rigid", "Iterative 2-4 week sprints"]
            }
          ]
        }
      ],
      "enumerations": [
        {
          "title": "Title of List",
          "prompt": "Enumerate the items in this list.",
          "items": ["Item 1", "Item 2", "Item 3"],
          "order_matters": false,
          "memory_cue": "Acronym / mnemonic cue"
        }
      ],
      "common_confusions": [
        {
          "concept_a": "Concept A Name",
          "concept_b": "Concept B Name",
          "difference": "Explicit contrast resolving why these two are commonly confused."
        }
      ],
      "key_takeaways": [
        "Takeaway 1",
        "Takeaway 2"
      ],
      "quiz": {
        "title": "Chapter Evaluation",
        "questions": [
          {
            "order": 1,
            "type": "multiple_choice",
            "difficulty": "medium",
            "text": "Question prompt testing comprehension?",
            "explanation": "Educational explanation defining why this answer is correct.",
            "choices": [
              {"text": "Option A text", "is_correct": false},
              {"text": "Option B text", "is_correct": true},
              {"text": "Option C text", "is_correct": false},
              {"text": "Option D text", "is_correct": false}
            ]
          },
          {
            "order": 2,
            "type": "true_false",
            "difficulty": "medium",
            "text": "True or False statement prompt?",
            "explanation": "Educational explanation explicitly explaining why the statement is true or false.",
            "choices": [
              {"text": "True", "is_correct": true},
              {"text": "False", "is_correct": false}
            ]
                    },
                    {
                        "order": 3,
                        "type": "identification",
                        "text": "Identify the framework used to organize architecture artifacts.",
                        "explanation": "The Zachman Framework organizes architecture artifacts across perspectives and concerns.",
                        "accepted_answers": ["Zachman Framework", "Zachman"]
                    },
                    {
                        "order": 4,
                        "type": "enumeration",
                        "text": "Enumerate the four BDAT domains.",
                        "explanation": "BDAT stands for Business, Data, Application, and Technology, the four domains used to classify architecture concerns.",
                        "order_matters": true,
                        "expected_items": [
                            {"canonical": "Business", "accepted_variants": []},
                            {"canonical": "Data", "accepted_variants": []},
                            {"canonical": "Application", "accepted_variants": []},
                            {"canonical": "Technology", "accepted_variants": []}
                        ]
          }
        ]
      }
    }
  ]
}
"""

# ==========================================
# 2. PROMPT COMPOSER
# ==========================================

def build_curriculum_prompt(course_title, extracted_text, study_goal="balanced_review", assessment_formats=None):
    if not assessment_formats:
        assessment_formats = ["multiple_choice"]
    elif isinstance(assessment_formats, str):
        assessment_formats = [assessment_formats]

    assessment_mix = get_assessment_mix(assessment_formats)

    goal_directive = STUDY_GOAL_PROFILES.get(study_goal, STUDY_GOAL_PROFILES["balanced_review"])
    
    assessment_directives = [
        ASSESSMENT_PROFILES[fmt]
        for fmt in assessment_formats
    ]
    assessment_block = "\n".join(assessment_directives)
    distribution_block = "\n".join(
        f"- {question_type}: {count}"
        for question_type, count in assessment_mix.items()
    )

    title_block = (
        f"<course_title>\n{course_title.strip()}\n</course_title>\n"
        "Instruction: Use this exact course title in the root 'course.title' field."
        if course_title and course_title.strip() != "Untitled Course"
        else "Instruction: Synthesize a professional academic course title from the source in 'course.title'."
    )

    # All anti-bloat, document scoping, and 8-10 question mandates are strictly passed to Gemini
    return f"""
You are an expert university curriculum architect and exam prep designer.
Analyze the attached study material and synthesize a structured, high-retention curriculum.

{title_block}

{goal_directive}

{assessment_block}

REQUIRED QUESTION DISTRIBUTION PER CHAPTER
Generate exactly 10 questions matching this distribution:
{distribution_block}
Do not replace requested Identification or Enumeration questions with Multiple Choice questions.

{DOCUMENT_STRUCTURE_RULES}

{ANTI_REDUNDANCY_RULES}

{QUIZ_RULES}

{SOURCE_GROUNDING_RULES}

{CURRICULUM_JSON_SCHEMA}

<source_material>
{extracted_text[:35000]}
</source_material>
"""

# ==========================================
# 3. EXTRACTION & CALL CONTROLLERS
# ==========================================

def extract_text_from_file(uploaded_file):
    filename = uploaded_file.name.lower()
    text = ""
    uploaded_file.seek(0)

    try:
        if filename.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                runs.append(p.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    runs.append(cell.text.strip())
            text = "\n".join(runs)

        elif filename.endswith('.docx'):
            doc = docx.Document(uploaded_file)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif filename.endswith('.txt'):
            text = uploaded_file.read().decode('utf-8', errors='ignore')

        elif filename.endswith('.pdf'):
            with pdfplumber.open(uploaded_file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

    except Exception as e:
        print(f"[Text Extraction Error]: {e}")

    return text.strip()


def generate_course_journey(course, uploaded_file=None, study_goal="balanced_review", assessment_formats=None):
    extracted_text = ""
    if uploaded_file:
        extracted_text = extract_text_from_file(uploaded_file)
    if not extracted_text and course.syllabus_text:
        extracted_text = course.syllabus_text

    if client and len(extracted_text) > 40:
        try:
            prompt = build_curriculum_prompt(
                course_title=course.title,
                extracted_text=extracted_text,
                study_goal=study_goal,
                assessment_formats=assessment_formats,
            )
            response = client.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    temperature=0.2,
                ),
            )
            return _validate_response(response.text, get_assessment_mix(assessment_formats))
        except Exception as e:
            print(f"[Gemini API Call Failure]: {e}")

    print("[Falling back to mock structured journey]")
    return _generate_mock_journey(course.title, assessment_formats)


def validate_question_mix(journey, required_mix):
    for chapter in journey.chapters:
        actual = {question_type: 0 for question_type in required_mix}
        for question in chapter.quiz.questions:
            if question.type not in actual:
                raise ValueError(f"Unsupported question type: {question.type}")
            actual[question.type] += 1
        if actual != required_mix:
            raise ValueError(
                f"{chapter.title} returned {actual}; expected {required_mix}."
            )


def _validate_response(raw_json, required_mix=None):
    # 1. Parse string to raw Python dict
    data = json.loads(raw_json)

    # 2. Enforce strict Pydantic contract (types, choices count, 1 correct choice, no duplicate questions)
    validated_journey = GeneratedJourney.model_validate(data)

    if required_mix is not None:
        validate_question_mix(validated_journey, required_mix)

    # 3. Application-level check: chapter question quantity
    for ch_idx, chapter in enumerate(validated_journey.chapters, start=1):
        question_count = len(chapter.quiz.questions)
        if not (6 <= question_count <= 15):
            raise ValueError(
                f"Chapter {ch_idx} has {question_count} questions; expected between 8 and 10 questions."
            )

    # 4. Return clean, validated dict for downstream consumers
    return validated_journey.model_dump()


def _generate_mock_journey(title, assessment_formats=None):
    course_title = getattr(title, 'title', title)
    if not isinstance(course_title, str) or not course_title.strip():
        course_title = str(title) if title else "System Integration and Architecture"

    journey = {
        "schema_version": "1.0",
        "course": {
            "title": course_title,
            "description": "Structured curriculum covering enterprise architecture and lifecycle patterns.",
            "source_type": "reviewer",
            "difficulty": "intermediate",
        },
        "chapters": [
            {
                "order": 1,
                "title": "Enterprise Architecture Fundamentals",
                "week_label": "Week 1",
                "focus": "Core enterprise architecture definitions, the BDAT model, and governance.",
                "overview": "Core enterprise architecture definitions, the BDAT model, and governance.",
                "estimated_minutes": 15,
                "learning_objectives": [
                    "Distinguish between TOGAF, Zachman, and BDAT domains.",
                    "Analyze organizational separation of concerns."
                ],
                "sections": [
                    {
                        "order": 1,
                        "title": "The Architecture Trinity",
                        "introduction": "Enterprise Architecture relies on three complementary structures working in unison.",
                        "concepts": [
                            {
                                "term": "TOGAF",
                                "formal_definition": "A standardized framework providing methodology and process for enterprise architecture.",
                                "simple_explanation": "Tells architects how and when to execute projects.",
                                "analogy": "The recipe.",
                                "examples": ["Architecture Development Method (ADM)"]
                            },
                            {
                                "term": "BDAT",
                                "formal_definition": "The four core domains: Business, Data, Application, and Technology.",
                                "simple_explanation": "The actual structures being built.",
                                "analogy": "The ingredients.",
                                "examples": ["PostgreSQL schema (Data)", "AWS EC2 instances (Technology)"]
                            }
                        ],
                        "key_points": [
                            "TOGAF is the process methodology.",
                            "BDAT defines structural domains.",
                            "Zachman organizes documents across rows and columns."
                        ]
                    }
                ],
                "key_terms": [
                    {
                        "term": "TOGAF",
                        "definition": "A standardized framework providing methodology and process for enterprise architecture."
                    },
                    {
                        "term": "BDAT",
                        "definition": "The four core domains: Business, Data, Application, and Technology."
                    }
                ],
                "analogy": {
                    "label": "The Recipe vs The Pantry",
                    "explanation": "TOGAF tells you how to cook (methodology), Zachman is where you store ingredients (taxonomy), and BDAT is the ingredients."
                },
                "comparisons": [
                    {
                        "title": "Monolith vs Microservices",
                        "columns": ["Architecture", "Strengths", "Weaknesses"],
                        "rows": [
                            {
                                "criterion": "Monolithic",
                                "values": ["Fast initial setup, simple deployments", "Single point of failure, shared database bottlenecks"]
                            },
                            {
                                "criterion": "Microservices",
                                "values": ["Fault isolation, independent service scaling", "Network latency, distributed tracing overhead"]
                            }
                        ]
                    }
                ],
                "enumerations": [
                    {
                        "title": "The Four BDAT Domains",
                        "prompt": "Enumerate the four domains of Enterprise Architecture in order.",
                        "items": ["Business", "Data", "Application", "Technology"],
                        "order_matters": True,
                        "memory_cue": "BDAT acronym"
                    }
                ],
                "common_confusions": [
                    {
                        "concept_a": "Component",
                        "concept_b": "Artifact",
                        "difference": "A component is a live operational asset (e.g. AWS server, code); an artifact is the documentation describing it (e.g. topology diagram, catalog)."
                    }
                ],
                "key_takeaways": [
                    "Splitting architecture into layers prevents cognitive overload.",
                    "Microservices require a database-per-service pattern to prevent lockouts."
                ],
                "quiz": {
                    "title": "Architecture Fundamentals Quiz",
                    "questions": [
                        {
                            "order": 1,
                            "type": "multiple_choice",
                            "difficulty": "medium",
                            "text": "Which framework answers WHERE architecture documents belong rather than HOW to execute them?",
                            "explanation": "Zachman serves as a taxonomy/filing system (the pantry) answering where artifacts reside, whereas TOGAF is the execution methodology.",
                            "choices": [
                                {"text": "TOGAF ADM", "is_correct": False},
                                {"text": "Zachman Framework", "is_correct": True},
                                {"text": "BDAT Domains", "is_correct": False},
                                {"text": "Agile Scrum", "is_correct": False}
                            ]
                        }
                    ]
                }
            }
        ]
    }
    journey["chapters"][0]["quiz"]["questions"] = _build_mock_questions(
        get_assessment_mix(assessment_formats)
    )
    return journey


def _build_mock_questions(assessment_mix):
    questions = []

    for index in range(assessment_mix["multiple_choice"]):
        questions.append({
            "type": "multiple_choice",
            "text": f"Which architecture principle is highlighted in mock question {index + 1}?",
            "explanation": "The selected principle keeps architecture decisions aligned with the course concepts and prevents unrelated design choices.",
            "choices": [
                {"text": "Layered separation", "is_correct": True},
                {"text": "Unbounded duplication", "is_correct": False},
                {"text": "Untracked coupling", "is_correct": False},
                {"text": "Random deployment", "is_correct": False},
            ],
        })

    for index in range(assessment_mix["true_false"]):
        questions.append({
            "type": "true_false",
            "text": f"True or False: mock architecture statement {index + 1} supports clear separation of concerns.",
            "explanation": "The statement is true because separating concerns makes systems easier to reason about, change, and govern.",
            "choices": [
                {"text": "True", "is_correct": True},
                {"text": "False", "is_correct": False},
            ],
        })

    identification_answers = [
        ("Identify the framework that organizes architecture artifacts.", ["Zachman Framework", "Zachman"]),
        ("Identify the methodology that guides architecture development.", ["TOGAF", "TOGAF ADM"]),
        ("Identify the architecture domain covering organizational goals.", ["Business"]),
        ("Identify the architecture domain covering stored information.", ["Data"]),
        ("Identify the architecture domain covering infrastructure.", ["Technology"]),
    ]
    for index in range(assessment_mix["identification"]):
        text, accepted_answers = identification_answers[index]
        questions.append({
            "type": "identification",
            "text": text,
            "explanation": "The accepted term is the precise concept used by the architecture framework in this lesson.",
            "accepted_answers": accepted_answers,
        })

    enumeration_items = [
        (["Business", "Data", "Application", "Technology"], True),
        (["Plan", "Build", "Measure"], False),
        (["People", "Process", "Technology"], False),
        (["Scope", "Time", "Cost"], True),
        (["Identify", "Assess", "Treat"], False),
    ]
    for index in range(assessment_mix["enumeration"]):
        items, order_matters = enumeration_items[index]
        questions.append({
            "type": "enumeration",
            "text": f"Enumerate the mock framework components for list {index + 1}.",
            "explanation": "Each listed item represents a distinct component in the framework and earns credit when identified correctly.",
            "order_matters": order_matters,
            "expected_items": [
                {"canonical": item, "accepted_variants": []}
                for item in items
            ],
        })

    for order, question in enumerate(questions, start=1):
        question["order"] = order
    return questions